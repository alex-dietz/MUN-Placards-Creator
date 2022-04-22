import copy
import random, os,uuid
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
import pandas as pd

pres = Presentation("placards.pptx")

def duplicate_slide(pres, index):
    template = pres.slides[index]
    try:
        blank_slide_layout = pres.slide_layouts[6]
    except:        
        blank_slide_layout = pres.slide_layouts[len(pres.slide_layouts)-1]

    copied_slide = pres.slides.add_slide(blank_slide_layout)
     # create images dict
    imgDict = {}

    # now copy contents from external slide, but do not copy slide properties
    # e.g. slide layouts, etc., because these would produce errors, as diplicate
    # entries might be generated
    for shp in template.shapes:
        
        if 'Logo' in shp.name:
            # save image

            with open(shp.name+'.jpg', 'wb') as f:
                
                f.write(shp.image.blob)

            # add image to dict
            imgDict[shp.name+'.jpg'] = [shp.left, shp.top, shp.width, shp.height]
        else:
            # create copy of elem
            el = shp.element
            newel = copy.deepcopy(el)

            # add elem to shape tree
            copied_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

    # add pictures
    for k, v in imgDict.items():
        copied_slide.shapes.add_picture(k, v[0], v[1], v[2], v[3])
        os.remove(k)

   
    return copied_slide

def find_shape_by_name(shapes, name):
    for shape in shapes:
        if shape.name == name:
            return shape
    return None





shapes = ['Committee','Country','Name']
delegates = pd.read_excel('delegates.xlsx')
print(len(delegates['Name']))

for index in range(len(delegates['Name'])):
    copy_slide = duplicate_slide(pres,0)
    for field in shapes:
        
        shape = find_shape_by_name(copy_slide.shapes,field)
        text_frame = shape.text_frame
        text_frame.clear() 
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = delegates[field][index]
        font = run.font
        font.name = 'SF Compact Display'
        font.size = Pt(14)
        font.italic = None  # cause value to be inherited from theme
        font.color.rgb = RGBColor(0, 0, 0)



pres.save('placards.pptx')